#!/usr/bin/env python3
"""Generate the Fluid 1.1.0 release presentation deck (.pptx) — English."""
import sys
sys.path.insert(0, '/Users/cheyang.cy/.claude/skills/ppt-architect/scripts')
from ppt_architect import *  # noqa: F401,F403
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm
from PIL import Image
import os

HERE = os.path.dirname(os.path.abspath(__file__))
NAVY = RGBColor(0x0B, 0x3D, 0x91)
SLIDE_W = 33.867  # cm (13.33")

ppt = PPTArchitect()
ppt.set_scale(0.9, 1.4, 1.6)  # logical canvas 34 x 18


def L(slide, x, y, w, h, text, size=11, bold=False, color=DARK):
    ppt.add_text(slide, x, y, w, h, text, size=size, bold=bold,
                 color=color, align=PP_ALIGN.LEFT)


def bg(slide, color=NAVY):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def pic(slide, name, w_cm, y_cm):
    path = os.path.join(HERE, name)
    iw, ih = Image.open(path).size
    x_cm = (SLIDE_W - w_cm) / 2
    slide.shapes.add_picture(path, Cm(x_cm), Cm(y_cm), width=Cm(w_cm))
    return w_cm * ih / iw


# ─────────────────────────────────────────────── Slide 1 · Title
sl = ppt.new_slide()
bg(sl)
ppt.add_text(sl, 2, 4.2, 30, 2.6, 'Fluid 1.1.0', size=54, bold=True, color=WHITE)
ppt.add_text(sl, 2, 7.2, 30, 1.6,
             'The CacheRuntime framework — onboard any cache system with zero controller code',
             size=18, bold=False, color=WHITE)
ppt.add_box(sl, 2, 10.2, 24, 1.8, RGBColor(0xFF, 0xC1, 0x07), YELLOW,
            '🎉  Jan 2026 · promoted to a CNCF Incubating project',
            13, RGBColor(0x1a, 0x1a, 0x1a), True)
ppt.add_text(sl, 2, 15.4, 30, 1.0,
             'CNCF Cloud Native Data Orchestration · github.com/fluid-cloudnative/fluid',
             size=12, color=RGBColor(0xC7, 0xD6, 0xF5))

# ─────────────────────────────────────────────── Slide 2 · At a glance
sl = ppt.new_slide('Release at a glance · since v1.0.8 (Oct 2025)')
stats = [
    ('485', 'Commits', LIGHT_BLUE, BLUE),
    ('487', 'Pull Requests', LIGHT_GREEN, GREEN),
    ('125', 'Contributors', LIGHT_ORANGE, ORANGE),
    ('~9', 'months', LIGHT_PURPLE, PURPLE),
]
bw, gap = 7.2, 1.4
x0 = (34 - (bw * 4 + gap * 3)) / 2
for i, (num, label, fill, line) in enumerate(stats):
    x = x0 + i * (bw + gap)
    ppt.add_box(sl, x, 4.5, bw, 4.2, fill, line, '', 10)
    ppt.add_text(sl, x, 5.0, bw, 2.0, num, size=40, bold=True, color=line)
    ppt.add_text(sl, x, 7.1, bw, 1.0, label, size=14, bold=True, color=DARK)
ppt.add_text(sl, 2, 10.0, 30, 1.2,
             'An architectural leap: from "write a controller per cache system" to '
             '"onboard any cache system with a single template".',
             size=14, bold=True, color=DARK)
ppt.add_text(sl, 2, 11.8, 30, 1.0, 'v1.0.8 (2025-10-30)  →  v1.1.0 (2026-07)',
             size=12, color=GRAY)

# ─────────────────────────────────────────────── Slide 3 · CacheRuntime topology
sl = ppt.new_slide('Flagship: the CacheRuntime generic cache-engine framework')
h = pic(sl, 'cacheruntime-arch-landscape-en.png', 31.5, 4.0)
ppt.add_text(sl, 2, 4.2 + h + 0.6, 30, 1.6,
             'CacheRuntimeClass (provider template) + CacheRuntime (user instance) + Dataset '
             '→ Controller → Master / Worker / Client.\n'
             'Onboarding an engine drops from "write a Go controller" to "fill in a YAML template".',
             size=13, bold=True, color=DARK)

# ─────────────────────────────────────────────── Slide 4 · capabilities
sl = ppt.new_slide('CacheRuntime · Core Capabilities')
caps = [
    ('Declarative & pluggable', 'Two new CRDs: CacheRuntimeClass (template) + CacheRuntime (instance) — no controller code', LIGHT_BLUE, BLUE),
    ('Built-in topologies', 'MasterSlave / P2P·DHT / ClientOnly — covers mainstream cache architectures', LIGHT_GREEN, GREEN),
    ('Tiered store', 'Both worker and client support Mem / SSD / HDD multi-tier caching', LIGHT_GREEN, GREEN),
    ('In-place upgrade', 'Via AdvancedStatefulSet — image changes roll in place, less cache rebuild', LIGHT_PURPLE, PURPLE),
    ('Data-flow ready', 'Natively wired to DataLoad / DataProcess for prefetch & processing', LIGHT_ORANGE, ORANGE),
    ('Dynamic mount & self-heal', 'Auto-remount on mount change / master restart; cache state synced to Dataset', LIGHT_YELLOW, YELLOW),
]
cw, ch, gx, gy = 15.2, 3.2, 1.4, 0.9
x0 = (34 - (cw * 2 + gx)) / 2
for i, (title, desc, fill, line) in enumerate(caps):
    r, c = divmod(i, 2)
    x = x0 + c * (cw + gx)
    y = 3.4 + r * (ch + gy)
    ppt.add_box(sl, x, y, cw, ch, fill, line, '', 9)
    L(sl, x + 0.5, y + 0.3, cw - 1, 0.9, '● ' + title, size=13, bold=True, color=line)
    L(sl, x + 0.5, y + 1.35, cw - 1, 1.6, desc, size=10.5, color=DARK)

# ─────────────────────────────────────────────── Slide 5 · Curvine
sl = ppt.new_slide('First engine: Curvine')
ppt.add_box(sl, 2, 3.4, 30, 2.2, LIGHT_GREEN, GREEN,
            'Curvine — a high-performance, Rust-based cloud-native distributed cache; '
            'the first engine to land on the CacheRuntime framework',
            13, DARK, True)
feats = [
    'Master / Worker architecture, Raft for metadata HA',
    'Multi-tier cache: memory / SSD / HDD',
    'POSIX-compatible access via FUSE',
    'Works with S3 / HDFS / OSS / MinIO backends',
    'Per community benchmarks: higher read throughput than open-source Alluxio on the same hardware',
    'A few lines of YAML spin up a full cache cluster via CacheRuntime',
]
for i, f in enumerate(feats):
    L(sl, 3, 6.4 + i * 1.25, 28, 1.1, '✓  ' + f, size=13, color=DARK)

# ─────────────────────────────────────────────── Slide 6 · AI acceleration
sl = ppt.new_slide('Data Acceleration for AI')
h = pic(sl, 'ai-data-pipeline-en.png', 31.5, 4.5)
ppt.add_text(sl, 2, 4.7 + h + 0.8, 30, 1.6,
             'Prefetch weights / training data into cache → cut inference cold-start, raise bandwidth use.\n'
             'High-performance stores (e.g. 3FS) integrate via ThinRuntime; KV-cache offload is a 2026 roadmap item.',
             size=13, bold=True, color=DARK)

# ─────────────────────────────────────────────── Slide 7 · Engine enhancements
sl = ppt.new_slide('Engine & Feature Enhancements')
items = [
    ('JindoCache', 'RDMA support\nunlocks high-perf network read\n+ master crash recovery', LIGHT_BLUE, BLUE),
    ('JuiceFSRuntime', 'Worker\nvolumeClaimTemplates\nmore flexible persistence', LIGHT_GREEN, GREEN),
    ('DataProcess', 'Cron scheduling\n+ OnEvent triggers\nfits automation pipelines', LIGHT_ORANGE, ORANGE),
]
bw, gap = 9.5, 1.6
x0 = (34 - (bw * 3 + gap * 2)) / 2
for i, (title, desc, fill, line) in enumerate(items):
    x = x0 + i * (bw + gap)
    ppt.add_box(sl, x, 4.5, bw, 5.0, fill, line, '', 9)
    ppt.add_text(sl, x, 5.0, bw, 1.0, title, size=15, bold=True, color=line)
    ppt.add_text(sl, x, 6.2, bw, 3.0, desc, size=12, color=DARK)
ppt.add_text(sl, 2, 10.8, 30, 1.4,
             'Also: global skip-runtime-sync switch · more resource names on runtime pods · '
             'host-PID feature gate · CSI fast node-label path & recovery loop · Jindo multi-mount OSS secrets',
             size=11, color=GRAY)

# ─────────────────────────────────────────────── Slide 8 · Reliability & Security
sl = ppt.new_slide('Reliability & Security Hardening')
ppt.add_box(sl, 1, 3.4, 15.5, 10, LIGHT_GREEN, GREEN, '', 9, DARK, False, True)
ppt.add_text(sl, 1, 3.6, 15.5, 0.9, 'Reliability & quality', size=14, bold=True, color=GREEN)
rel = [
    '240+ testing / E2E improvements (largest ever)',
    'Full test-suite migration to Ginkgo v2 + Gomega',
    'gohook → gomonkey for stubbing',
    'Fixed data races and controller panics',
    'Retry-on-conflict on critical paths (UFS / webhook)',
    'DataLoad/DataMigrate correctly watch Jobs',
]
for i, t in enumerate(rel):
    L(sl, 1.6, 4.8 + i * 1.3, 14.5, 1.1, '• ' + t, size=11, color=DARK)
ppt.add_box(sl, 17.5, 3.4, 15.5, 10, LIGHT_RED, RED, '', 9, DARK, False, True)
ppt.add_text(sl, 17.5, 3.6, 15.5, 0.9, 'Security hardening', size=14, bold=True, color=RED)
sec = [
    'Least-privilege RBAC / ClusterRole',
    'Tightened GitHub Actions permissions',
    'Fixed path-injection risk in addons',
    'Base image bumped to alpine 3.23.3',
    'Removed bundled Helm binaries from repo',
    'Removed hardcoded Asia/Shanghai timezone',
]
for i, t in enumerate(sec):
    L(sl, 18.1, 4.8 + i * 1.3, 14.5, 1.1, '• ' + t, size=11, color=DARK)

# ─────────────────────────────────────────────── Slide 9 · Deprecation & upgrade
sl = ppt.new_slide('Deprecation & Upgrade')
ppt.add_box(sl, 2, 3.6, 30, 2.2, LIGHT_YELLOW, YELLOW,
            '⚠  GooseFSRuntime is deprecated; its resources and docs are removed — plan migration',
            13, DARK, True)
ppt.add_text(sl, 2, 6.6, 30, 0.9, 'Upgrade:', size=14, bold=True, color=DARK)
ppt.add_box(sl, 2, 7.5, 30, 3.6, RGBColor(0x2B, 0x2B, 0x2B), RGBColor(0x2B, 0x2B, 0x2B),
            'helm repo add fluid https://fluid-cloudnative.github.io/charts\n'
            'helm repo update\n'
            'helm upgrade --install fluid fluid/fluid --version 1.1.0 -n fluid-system',
            13, WHITE, False)
ppt.add_text(sl, 2, 11.6, 30, 1.6,
             'Note: new CacheRuntime / CacheRuntimeClass CRDs — update CRDs on upgrade;\n'
             'using CacheRuntime requires the AdvancedStatefulSet (OpenKruise) controller.',
             size=12, color=GRAY)

# ─────────────────────────────────────────────── Slide 10 · Roadmap
sl = ppt.new_slide('2026 Roadmap')
ppt.add_box(sl, 1, 3.6, 11.5, 4.6, LIGHT_GREEN, GREEN, '', 9, DARK, False, True)
ppt.add_text(sl, 1.3, 3.7, 11, 0.9, 'Data Anyway · regardless of infra', size=11.5, bold=True, color=GREEN)
for i, t in enumerate([
        'Generic Cache Runtime pluggable interface',
        'Full migration to AdvancedStatefulSet',
        'Zero-downtime runtime tuning',
        'API upgrade to v1alpha2 + webhook',
        'ThinRuntime productization (min privilege)']):
    L(sl, 1.6, 4.7 + i * 0.66, 10.8, 0.6, '▸ ' + t, size=9.5, color=DARK)
ppt.add_box(sl, 1, 8.8, 11.5, 4.2, LIGHT_BLUE, BLUE, '', 9, DARK, False, True)
ppt.add_text(sl, 1.3, 8.9, 11, 0.9, 'Data Anywhere · cross-region/cluster', size=11.5, bold=True, color=BLUE)
for i, t in enumerate([
        'LLM KV Cache orchestration (vLLM/SGLang · Mooncake)',
        'Cross-pod KV cache migration & sharing',
        'Distributed prewarming + throttling',
        'Cross-region sync optimization']):
    L(sl, 1.6, 9.9 + i * 0.72, 10.8, 0.6, '▸ ' + t, size=9.5, color=DARK)
path = os.path.join(HERE, 'roadmap-2026-en.png')
iw, ih = Image.open(path).size
img_h = 9.6
img_w = img_h * iw / ih
sl.shapes.add_picture(path, Cm(13.5 + (20 - img_w) / 2), Cm(3.6), height=Cm(img_h))

# ─────────────────────────────────────────────── Slide 11 · Get involved
sl = ppt.new_slide()
bg(sl)
ppt.add_text(sl, 2, 4.5, 30, 2.0, 'Get Involved', size=40, bold=True, color=WHITE)
ppt.add_text(sl, 2, 7.4, 30, 3.5,
             '⭐  github.com/fluid-cloudnative/fluid\n\n'
             '💬  Open issues / PRs — bring your cache engine to the CacheRuntime ecosystem\n\n'
             '🙏  Thanks to the 125 contributors of this cycle',
             size=17, color=WHITE)
ppt.add_text(sl, 2, 14.5, 30, 1.2, 'Fluid 1.1.0 · CNCF Incubating Project',
             size=14, bold=True, color=RGBColor(0xFF, 0xC1, 0x07))

ppt.save('docs/blog/v1.1.0/fluid-1.1.0-deck-en.pptx')
print('saved')
