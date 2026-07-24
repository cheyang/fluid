#!/usr/bin/env python3
"""Generate the Fluid 1.1.0 release presentation deck (.pptx)."""
import sys
sys.path.insert(0, '/Users/cheyang.cy/.claude/skills/ppt-architect/scripts')
from ppt_architect import *  # noqa: F401,F403
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm
from PIL import Image
import os

HERE = os.path.dirname(os.path.abspath(__file__))
NAVY = RGBColor(0x0B, 0x3D, 0x91)
SLIDE_W = 33.867  # cm (13.33")

ppt = PPTArchitect()
ppt.set_scale(0.9, 1.4, 1.6)  # logical canvas 34 x 18


def L(slide, x, y, w, h, text, size=11, bold=False, color=DARK):
    ppt.add_text(slide, x, y, w, h, text, size=size, bold=bold,
                 color=color, align=PP_ALIGN.LEFT)


def bg(slide, color=NAVY):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def pic(slide, name, w_cm, y_cm):
    """Center a picture horizontally at width w_cm; return its height in cm."""
    path = os.path.join(HERE, name)
    iw, ih = Image.open(path).size
    x_cm = (SLIDE_W - w_cm) / 2
    slide.shapes.add_picture(path, Cm(x_cm), Cm(y_cm), width=Cm(w_cm))
    return w_cm * ih / iw


# ─────────────────────────────────────────────── Slide 1 · Title
sl = ppt.new_slide()
bg(sl)
ppt.add_text(sl, 2, 4.2, 30, 2.6, 'Fluid 1.1.0', size=54, bold=True, color=WHITE)
ppt.add_text(sl, 2, 7.2, 30, 1.6,
             'CacheRuntime 通用缓存引擎框架 · 让任意缓存系统「零代码」接入',
             size=20, bold=False, color=WHITE)
ppt.add_box(sl, 2, 10.2, 22, 1.8, RGBColor(0xFF, 0xC1, 0x07), YELLOW,
            '🎉 2026.01 · 正式晋升为 CNCF 孵化 (Incubating) 项目',
            13, RGBColor(0x1a, 0x1a, 0x1a), True)
ppt.add_text(sl, 2, 15.4, 30, 1.0,
             'CNCF Cloud Native Data Orchestration · github.com/fluid-cloudnative/fluid',
             size=12, color=RGBColor(0xC7, 0xD6, 0xF5))

# ─────────────────────────────────────────────── Slide 2 · At a glance
sl = ppt.new_slide('版本速览 · 自 v1.0.8 (2025.10) 起')
stats = [
    ('485', 'Commits', LIGHT_BLUE, BLUE),
    ('487', 'Pull Requests', LIGHT_GREEN, GREEN),
    ('125', 'Contributors', LIGHT_ORANGE, ORANGE),
    ('~9', '个月迭代', LIGHT_PURPLE, PURPLE),
]
bw, gap = 7.2, 1.4
x0 = (34 - (bw * 4 + gap * 3)) / 2
for i, (num, label, fill, line) in enumerate(stats):
    x = x0 + i * (bw + gap)
    ppt.add_box(sl, x, 4.5, bw, 4.2, fill, line, '', 10)
    ppt.add_text(sl, x, 5.0, bw, 2.0, num, size=40, bold=True, color=line)
    ppt.add_text(sl, x, 7.1, bw, 1.0, label, size=14, bold=True, color=DARK)
ppt.add_text(sl, 2, 10.0, 30, 1.2,
             '一次架构级跃迁：从「为每种缓存系统写一个控制器」到「用一份模板接入任意缓存系统」',
             size=15, bold=True, color=DARK)
ppt.add_text(sl, 2, 11.6, 30, 1.0,
             'v1.0.8 (2025-10-30)  →  v1.1.0 (2026-07)',
             size=12, color=GRAY)

# ─────────────────────────────────────────────── Slide 3 · CacheRuntime topology
sl = ppt.new_slide('旗舰特性：CacheRuntime 通用缓存引擎框架')
h = pic(sl, 'cacheruntime-arch-landscape.png', 31.5, 4.0)
ppt.add_text(sl, 2, 4.2 + h + 0.6, 30, 1.6,
             'CacheRuntimeClass(引擎方模板) + CacheRuntime(使用方实例) + Dataset → Controller '
             '→ Master / Worker / Client。\n把「接入一个缓存引擎」从「写一个 Go 控制器」降维成「填一份 YAML 模板」。',
             size=13, bold=True, color=DARK)

# ─────────────────────────────────────────────── Slide 4 · CacheRuntime capabilities
sl = ppt.new_slide('CacheRuntime · 核心能力')
caps = [
    ('声明式可插拔', '两个新 CRD：CacheRuntimeClass(模板) + CacheRuntime(实例)，无需编写控制器', LIGHT_BLUE, BLUE),
    ('多拓扑内建', 'MasterSlave / P2P·DHT / ClientOnly，覆盖主流缓存架构', LIGHT_GREEN, GREEN),
    ('分层存储', 'Worker 与 Client 均支持 Mem / SSD / HDD 多级缓存', LIGHT_GREEN, GREEN),
    ('原地升级', '基于 AdvancedStatefulSet，镜像变更原地滚动，减少缓存重建', LIGHT_PURPLE, PURPLE),
    ('数据流打通', '原生对接 DataLoad / DataProcess 预热与处理', LIGHT_ORANGE, ORANGE),
    ('动态挂载与自愈', '挂载点变更/Master 重启自动重挂；缓存状态实时回填 Dataset', LIGHT_YELLOW, YELLOW),
]
cw, ch, gx, gy = 15.2, 3.2, 1.4, 0.9
x0 = (34 - (cw * 2 + gx)) / 2
for i, (title, desc, fill, line) in enumerate(caps):
    r, c = divmod(i, 2)
    x = x0 + c * (cw + gx)
    y = 3.4 + r * (ch + gy)
    ppt.add_box(sl, x, y, cw, ch, fill, line, '', 9)
    L(sl, x + 0.5, y + 0.3, cw - 1, 0.9, '● ' + title, size=13, bold=True, color=line)
    L(sl, x + 0.5, y + 1.35, cw - 1, 1.6, desc, size=10.5, color=DARK)

# ─────────────────────────────────────────────── Slide 5 · Curvine
sl = ppt.new_slide('首发引擎：Curvine')
ppt.add_box(sl, 2, 3.4, 30, 2.2, LIGHT_GREEN, GREEN,
            'Curvine — 基于 Rust 的高性能云原生分布式缓存系统，CacheRuntime 框架的首个完整落地引擎',
            14, DARK, True)
feats = [
    'Master / Worker 架构，Raft 保障元数据高可用',
    '多级缓存：内存 / SSD / HDD',
    'FUSE 提供 POSIX 兼容访问',
    '兼容 S3 / HDFS / OSS / MinIO 等底层存储',
    '据社区基准：相同硬件下读吞吐优于开源版 Alluxio',
    '通过 CacheRuntime，几行 YAML 即可拉起一套缓存集群',
]
for i, f in enumerate(feats):
    L(sl, 3, 6.4 + i * 1.25, 28, 1.1, '✓  ' + f, size=13, color=DARK)

# ─────────────────────────────────────────────── Slide 6 · AI acceleration
sl = ppt.new_slide('面向 AI 场景的数据加速')
h = pic(sl, 'ai-data-pipeline-zh.png', 31.5, 4.5)
ppt.add_text(sl, 2, 4.7 + h + 0.8, 30, 1.6,
             '模型权重 / 训练数据预热进缓存 → 缩短推理冷启动、提升带宽利用率；\n'
             '高性能存储(3FS 等)经 ThinRuntime 接入；KV Cache 外置为 2026 Roadmap 方向。',
             size=13, bold=True, color=DARK)

# ─────────────────────────────────────────────── Slide 7 · Engine enhancements
sl = ppt.new_slide('引擎与功能增强')
items = [
    ('JindoCache', 'RDMA 支持\n释放高性能网络读带宽\n+ Master 崩溃自愈', LIGHT_BLUE, BLUE),
    ('JuiceFSRuntime', 'Worker\nvolumeClaimTemplates\n更灵活的持久化存储', LIGHT_GREEN, GREEN),
    ('DataProcess', 'Cron 定时触发\n+ OnEvent 事件触发\n融入自动化运维', LIGHT_ORANGE, ORANGE),
]
bw, gap = 9.5, 1.6
x0 = (34 - (bw * 3 + gap * 2)) / 2
for i, (title, desc, fill, line) in enumerate(items):
    x = x0 + i * (bw + gap)
    ppt.add_box(sl, x, 4.5, bw, 5.0, fill, line, '', 9)
    ppt.add_text(sl, x, 5.0, bw, 1.0, title, size=15, bold=True, color=line)
    ppt.add_text(sl, x, 6.2, bw, 3.0, desc, size=12, color=DARK)
ppt.add_text(sl, 2, 10.8, 30, 1.4,
             '此外：全局跳过 Runtime 同步开关 · Runtime Pod 支持更多资源名 · host PID 特性门控 · CSI 快取节点标签与恢复循环 · Jindo 多挂载 OSS Secret',
             size=11, color=GRAY)

# ─────────────────────────────────────────────── Slide 8 · Reliability & Security
sl = ppt.new_slide('稳定性与安全加固')
ppt.add_box(sl, 1, 3.4, 15.5, 10, LIGHT_GREEN, GREEN, '', 9, DARK, False, True)
ppt.add_text(sl, 1, 3.6, 15.5, 0.9, '稳定性与工程质量', size=14, bold=True, color=GREEN)
rel = [
    '240+ 项测试与 E2E 改进（历史最大投入）',
    '测试套件全面迁移 Ginkgo v2 + Gomega',
    'gohook → gomonkey 打桩',
    '修复数据竞争、控制器 panic',
    '关键路径冲突重试（UFS 状态 / Webhook）',
    'DataLoad/DataMigrate 正确 watch Job',
]
for i, t in enumerate(rel):
    L(sl, 1.6, 4.8 + i * 1.3, 14.5, 1.1, '• ' + t, size=11.5, color=DARK)
ppt.add_box(sl, 17.5, 3.4, 15.5, 10, LIGHT_RED, RED, '', 9, DARK, False, True)
ppt.add_text(sl, 17.5, 3.6, 15.5, 0.9, '安全加固', size=14, bold=True, color=RED)
sec = [
    '控制面最小权限 RBAC / ClusterRole',
    'GitHub Actions 权限全面收敛',
    '修复 addons 路径注入风险',
    '基础镜像升级至 alpine 3.23.3',
    '移除仓库内置 Helm 二进制',
    '移除硬编码 Asia/Shanghai 时区',
]
for i, t in enumerate(sec):
    L(sl, 18.1, 4.8 + i * 1.3, 14.5, 1.1, '• ' + t, size=11.5, color=DARK)

# ─────────────────────────────────────────────── Slide 9 · Deprecation & upgrade
sl = ppt.new_slide('弃用说明与升级')
ppt.add_box(sl, 2, 3.6, 30, 2.2, LIGHT_YELLOW, YELLOW,
            '⚠ GooseFSRuntime 已废弃 (deprecated)，相关资源与文档已清理，建议存量用户规划迁移',
            13, DARK, True)
ppt.add_text(sl, 2, 6.6, 30, 0.9, '升级：', size=14, bold=True, color=DARK)
ppt.add_box(sl, 2, 7.5, 30, 3.6, RGBColor(0x2B, 0x2B, 0x2B), RGBColor(0x2B, 0x2B, 0x2B),
            'helm repo add fluid https://fluid-cloudnative.github.io/charts\n'
            'helm repo update\n'
            'helm upgrade --install fluid fluid/fluid --version 1.1.0 -n fluid-system',
            13, WHITE, False)
ppt.add_text(sl, 2, 11.6, 30, 1.6,
             '注意：新增 CacheRuntime / CacheRuntimeClass CRD，升级时请一并更新 CRD；\n'
             '使用 CacheRuntime 需在集群安装 AdvancedStatefulSet (OpenKruise) 控制器。',
             size=12, color=GRAY)

# ─────────────────────────────────────────────── Slide 10 · Roadmap
sl = ppt.new_slide('2026 Roadmap')
# left: two-branch bullets; right: rendered roadmap tree
ppt.add_box(sl, 1, 3.6, 11.5, 4.6, LIGHT_GREEN, GREEN, '', 9, DARK, False, True)
ppt.add_text(sl, 1.3, 3.7, 11, 0.9, 'Data Anyway · 不受基础设施约束', size=11.5, bold=True, color=GREEN)
for i, t in enumerate([
        'Generic Cache Runtime 可插拔接口',
        '全面迁移 AdvancedStatefulSet',
        'Runtime 零停机动态调参',
        'API 升级 v1alpha2 + 转换 Webhook',
        'ThinRuntime 产品化 (最小权限)']):
    L(sl, 1.6, 4.7 + i * 0.66, 10.8, 0.6, '▸ ' + t, size=9.5, color=DARK)
ppt.add_box(sl, 1, 8.8, 11.5, 4.2, LIGHT_BLUE, BLUE, '', 9, DARK, False, True)
ppt.add_text(sl, 1.3, 8.9, 11, 0.9, 'Data Anywhere · 跨地域/跨集群', size=11.5, bold=True, color=BLUE)
for i, t in enumerate([
        'LLM KV Cache 编排 (vLLM/SGLang · Mooncake)',
        '跨 Pod KV Cache 迁移共享',
        '分布式预热 + 带宽限速',
        '跨地域同步优化']):
    L(sl, 1.6, 9.9 + i * 0.72, 10.8, 0.6, '▸ ' + t, size=9.5, color=DARK)
# right: rendered tree image
path = os.path.join(HERE, 'roadmap-2026-zh.png')
iw, ih = Image.open(path).size
img_h = 9.6
img_w = img_h * iw / ih
sl.shapes.add_picture(path, Cm(13.5 + (20 - img_w) / 2), Cm(3.6), height=Cm(img_h))

# ─────────────────────────────────────────────── Slide 11 · Get involved
sl = ppt.new_slide()
bg(sl)
ppt.add_text(sl, 2, 4.5, 30, 2.0, '参与社区', size=40, bold=True, color=WHITE)
ppt.add_text(sl, 2, 7.4, 30, 3.5,
             '⭐  github.com/fluid-cloudnative/fluid\n\n'
             '💬  提交 Issue / PR，把你的缓存引擎接入 CacheRuntime 生态\n\n'
             '🙏  感谢本周期 125 位贡献者的付出',
             size=17, color=WHITE)
ppt.add_text(sl, 2, 14.5, 30, 1.2, 'Fluid 1.1.0 · CNCF Incubating Project',
             size=14, bold=True, color=RGBColor(0xFF, 0xC1, 0x07))

ppt.save('docs/blog/v1.1.0/fluid-1.1.0-deck.pptx')
print('saved')
